#!/usr/bin/env python3
"""펜 잡기 Sim2Real 프로젝트 발표자료 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 색상 정의
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 112, 192)
DARK_GRAY = RGBColor(64, 64, 64)
WHITE = RGBColor(255, 255, 255)
ACCENT_ORANGE = RGBColor(255, 153, 0)


def add_title_slide(prs, title, subtitle=""):
    """제목 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # 배경색 (진한 파란색)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 부제목
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, content_lines, bullet=True):
    """내용 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # 상단 바
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = DARK_BLUE
    top_bar.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 내용
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # 들여쓰기 처리
        if line.startswith("  - "):
            p.text = line.strip()[2:]
            p.level = 1
        elif line.startswith("- "):
            p.text = line.strip()[2:]
            p.level = 0
        else:
            p.text = line
            p.level = 0

        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)

    return slide


def add_table_slide(prs, title, headers, rows):
    """테이블 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # 상단 바
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = DARK_BLUE
    top_bar.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 테이블
    cols = len(headers)
    table_rows = len(rows) + 1  # +1 for header

    table_width = Inches(9)
    table_height = Inches(0.5 * table_rows)
    left = Inches(0.5)
    top = Inches(1.8)

    table = slide.shapes.add_table(table_rows, cols, left, top, table_width, table_height).table

    # 헤더 설정
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # 데이터 행
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER

    return slide


def add_diagram_slide(prs, title, diagram_text):
    """다이어그램/코드 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # 상단 바
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = DARK_BLUE
    top_bar.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 다이어그램 박스
    diagram_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    )
    diagram_box.fill.solid()
    diagram_box.fill.fore_color.rgb = RGBColor(245, 245, 245)

    # 다이어그램 텍스트
    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(8.6), Inches(4.6))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = diagram_text
    p.font.size = Pt(14)
    p.font.name = "Consolas"
    p.font.color.rgb = DARK_GRAY

    return slide


def create_presentation():
    """전체 프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. 표지
    add_title_slide(
        prs,
        "펜 잡기 Sim2Real 프로젝트",
        "로봇프로젝트 중간발표 | 2025.12.29"
    )

    # 2. 목차
    add_content_slide(prs, "목차", [
        "1. 프로젝트 개요",
        "2. 시스템 아키텍처",
        "3. Phase 1: 로봇 시스템 구축 (12/9~12/12)",
        "4. Phase 2: 시뮬레이션 환경 구축 (12/12~12/17)",
        "5. Phase 3: 강화학습 환경 발전 (12/17~12/24)",
        "6. Phase 4: Sim2Real 파이프라인 (12/22~12/27)",
        "7. 주요 기술적 도전과 해결",
        "8. 현재 결과 및 향후 계획",
    ])

    # 3. 프로젝트 개요
    add_content_slide(prs, "1. 프로젝트 개요", [
        "목표: 로봇이 책상 위의 펜을 인식하고 잡아서 글쓰기 자세로 만들기",
        "",
        "핵심 기술:",
        "- 강화학습 (RL): Isaac Lab 기반 시뮬레이션 학습",
        "- Sim2Real Transfer: 시뮬레이션 → 실제 로봇 정책 전이",
        "- 비전 시스템: YOLO + RealSense 기반 펜 감지",
        "",
        "하드웨어:",
        "- 로봇팔: Doosan E0509 (6축)",
        "- 그리퍼: Robotis RH-P12-RN-A",
        "- 카메라: Intel RealSense D455F",
    ])

    # 4. 시스템 아키텍처
    add_diagram_slide(prs, "2. 시스템 아키텍처 (3개 레포 구조)", """
┌─────────────────────────────────────────────────────────────────────┐
│                      펜 잡기 Sim2Real 프로젝트                        │
├─────────────────────┬─────────────────────┬─────────────────────────┤
│  e0509_gripper_     │    CoWriteBotRL     │        sim2real         │
│    description      │                     │                         │
├─────────────────────┼─────────────────────┼─────────────────────────┤
│ • ROS2 로봇 패키지   │ • Isaac Lab RL 환경  │ • YOLO 펜 감지           │
│ • URDF/XACRO        │ • 강화학습 환경       │ • Hand-Eye Calibration  │
│ • 그리퍼 제어        │ • 학습 스크립트       │ • Jacobian IK           │
│ • Digital Twin      │ • V1~V7 환경 진화    │ • 정책 실행 파이프라인    │
├─────────────────────┼─────────────────────┼─────────────────────────┤
│    ROS2 Humble      │   Isaac Sim 4.5     │    ROS2 + PyTorch       │
└─────────────────────┴─────────────────────┴─────────────────────────┘
""")

    # 5. Phase 1
    add_table_slide(prs, "3. Phase 1: 로봇 시스템 구축 (12/9~12/12)",
        ["날짜", "작업", "설명"],
        [
            ["12/9", "초기 통합", "E0509 + RH-P12-RN-A 그리퍼 URDF 결합"],
            ["12/9", "Gazebo 지원", "Ignition Gazebo 6 시뮬레이션 추가"],
            ["12/10", "Docker 설치", "Doosan Virtual Mode 에뮬레이터"],
            ["12/12", "실제 그리퍼 제어", "Tool Flange Serial + Modbus RTU"],
        ]
    )

    # 6. Phase 1 상세
    add_content_slide(prs, "Phase 1: 기술적 성과", [
        "ros2_control 기반 조인트 제어:",
        "- Doosan 로봇 SDK와 ROS2 통합",
        "- Virtual Mode와 Real Mode 지원",
        "",
        "그리퍼 통신 구현:",
        "- Tool Flange Serial 포트 활용",
        "- Modbus RTU 프로토콜 직접 구현",
        "- Baudrate: 57600, 8N1",
        "- Stroke 제어: 0(열림) ~ 700(닫힘)",
        "",
        "ROS2 서비스 인터페이스:",
        "- /gripper/open, /gripper/close 서비스",
        "- /gripper/position_cmd 토픽",
    ])

    # 7. Phase 2
    add_table_slide(prs, "4. Phase 2: 시뮬레이션 환경 구축 (12/12~12/17)",
        ["날짜", "작업", "설명"],
        [
            ["12/12", "Isaac Lab 환경", "pen_grasp_rl 패키지 생성"],
            ["12/15", "Digital Twin", "실제 로봇 → Isaac Sim 동기화"],
            ["12/17", "Direct 환경", "단계별 상태 머신 구현"],
            ["12/17", "보상 설계", "거리/정렬/진행 보상 구조"],
        ]
    )

    # 8. Phase 2 상세
    add_diagram_slide(prs, "Phase 2: Direct 환경 상태 머신", """
초기 버전 (V1):

    ┌──────────┐      ┌──────────┐      ┌──────────┐      ┌──────────┐
    │ APPROACH │  →   │  ALIGN   │  →   │  GRASP   │  →   │ SUCCESS! │
    │  (접근)   │      │  (정렬)   │      │  (잡기)   │      │          │
    └──────────┘      └──────────┘      └──────────┘      └──────────┘
         │                 │                 │
         ▼                 ▼                 ▼
    거리 < 10cm       dot < -0.8      거리 < 2cm &
                                      dot < -0.9


USD 모델:
  • e0509_gripper_isaac.usd (로봇)
  • pen.usd (펜 모델)

학습 설정:
  • PPO (RSL-RL)
  • 4096 병렬 환경
  • Actor/Critic: [256, 256, 128]
""")

    # 9. Phase 3 개요
    add_content_slide(prs, "5. Phase 3: 강화학습 환경 발전 (12/17~12/24)", [
        "핵심 전환: Joint Space → Task Space 제어",
        "",
        "문제점 발견:",
        "- Joint Space 제어의 한계",
        "- 관절 공간에서 학습 → Sim2Real 전이 어려움",
        "- 로봇 기구학적 특성에 과적합",
        "",
        "해결책: IK 환경 도입 (12/18)",
        "- Task Space (3DoF 위치) 제어로 전환",
        "- RL이 TCP 위치 변화량(delta)을 출력",
        "- Inverse Kinematics로 관절 각도 계산",
        "- 로봇에 독립적인 정책 학습 가능",
    ])

    # 10. IK 버전 진화
    add_table_slide(prs, "Phase 3: IK 환경 버전 진화",
        ["버전", "날짜", "핵심 변경"],
        [
            ["IK V1", "12/18", "Task Space Control 도입"],
            ["IK V2", "12/18", "ALIGN 단계 추가, 펜 캡 위치 기준"],
            ["IK V3", "12/19", "펜 축 기준 접근, 충돌 감지"],
            ["IK V4", "12/19", "Hybrid RL + TCP, 펜 각도 랜덤화 (±30°)"],
            ["IK V5", "12/22", "TCP 버그 수정, Curriculum Learning"],
            ["IK V6", "12/23", "3DoF 위치 제어 + 자동 자세 정렬"],
            ["IK V7", "12/24", "APPROACH Only - Sim2Real Ready"],
        ]
    )

    # 11. IK V7 상세
    add_diagram_slide(prs, "IK V7: Sim2Real 최적화 설계", """
핵심 철학: "Simple is Best"

┌─────────────────────────────────────────────────────────────────────┐
│  기존 (V1~V6):                                                       │
│  APPROACH → ALIGN → DESCEND → GRASP → LIFT → SUCCESS                │
│                              ↓ 간소화                                │
│  V7:                                                                 │
│  APPROACH (접근만) → 그리퍼 닫기 (조건부, 하드코딩)                    │
└─────────────────────────────────────────────────────────────────────┘

왜 간소화?
  1. Sim2Real Gap 최소화: 복잡한 상태 머신 → 전이 어려움
  2. 그리퍼 동작 분리: RL은 접근만, 그리퍼는 거리 조건부 닫기
  3. 학습 안정성: 단순한 목표 → 빠른 수렴

관찰 공간 (15차원):
  • TCP 위치 (3D) + 펜 캡 위치 (3D) + 상대 위치 (3D)
  • 그리퍼 Z축 (3D) + 펜 Z축 (3D)

액션 공간 (3차원):
  • TCP delta position (x, y, z)
""")

    # 12. Phase 4
    add_diagram_slide(prs, "6. Phase 4: Sim2Real 파이프라인 (12/22~12/27)", """
┌──────────────┐    ┌──────────────┐    ┌──────────────┐
│   RealSense  │ →  │  YOLO 펜감지  │ →  │   좌표 변환   │
│    D455F     │    │ Segmentation │    │  Cam→Robot   │
└──────────────┘    └──────────────┘    └──────────────┘
                                               │
                                               ▼
┌──────────────┐    ┌──────────────┐    ┌──────────────┐
│  그리퍼 닫기  │ ←  │ Jacobian IK  │ ←  │   RL 정책    │
│              │    │  (DLS 방식)   │    │  (PyTorch)   │
└──────────────┘    └──────────────┘    └──────────────┘


핵심 컴포넌트:
  1. YOLO 펜 감지: YOLOv8 Segmentation → 펜 캡/팁 위치 추출
  2. Hand-Eye Calibration: Eye-to-Hand, ArUco 마커 기반
  3. Jacobian IK: Damped Least Squares (DLS) 방식
""")

    # 13. 기술적 도전
    add_table_slide(prs, "7. 주요 기술적 도전과 해결",
        ["도전", "문제", "해결"],
        [
            ["그리퍼 통신", "Doosan SDK에 그리퍼 제어 없음", "Tool Flange Serial + Modbus RTU 직접 구현"],
            ["Digital Twin", "ROS2 (Py3.10) ↔ Isaac (Py3.11)", "파일 기반 통신 (JSON)"],
            ["Sim2Real Gap", "복잡한 상태 머신 전이 실패", "V7: APPROACH Only 간소화"],
            ["Joint 과적합", "관절 공간 학습의 한계", "Task Space (IK) 제어 전환"],
        ]
    )

    # 14. 현재 결과
    add_content_slide(prs, "8. 현재 결과", [
        "학습 결과 (V7, 100k iterations):",
        "- 병렬 환경 수: 4096",
        "- 학습 시간: ~6시간",
        "- Mean Reward: 수렴",
        "- 성공률 (시뮬레이션): ~90%+",
        "",
        "Sim2Real 테스트 현황:",
        "- YOLO 펜 감지: 성공",
        "- Hand-Eye Calibration: 성공",
        "- 좌표 변환: 성공",
        "- 로봇 접근 동작: 진행 중",
    ])

    # 15. 향후 계획
    add_content_slide(prs, "향후 계획", [
        "단기 (1~2주):",
        "- Sim2Real 테스트 완료: 실제 로봇에서 펜 잡기 성공",
        "- 그리퍼 타이밍 최적화: 접근 거리 기반 그리퍼 닫기",
        "",
        "중기 (1달):",
        "- 다양한 펜 위치/각도: Domain Randomization 강화",
        "- 글쓰기 동작 추가: 펜 잡기 → 글쓰기 자세 전환",
        "",
        "장기:",
        "- End-to-End 비전 정책: 카메라 이미지 → 동작 직접 학습",
        "- 실제 글쓰기 작업: 획 단위 제어",
    ])

    # 16. 타임라인 요약
    add_diagram_slide(prs, "프로젝트 타임라인 요약", """
12/9 ─────────── 12/12 ─────────── 12/17 ─────────── 12/24 ─────────── 12/27
  │                │                 │                 │                 │
  ▼                ▼                 ▼                 ▼                 ▼
로봇 시스템       그리퍼 제어        Isaac Lab        IK V7           Sim2Real
  구축              완성            환경 구축          완성           파이프라인


핵심 성과:

  1. 3개 레포 아키텍처
     → 역할별 분리로 효율적 개발

  2. IK 기반 제어
     → Sim2Real 전이 최적화

  3. V7 환경
     → 단순화된 접근 전략으로 안정적 학습

  4. 완전한 파이프라인
     → 감지 → 정책 → IK → 제어
""")

    # 17. Q&A
    add_title_slide(
        prs,
        "Q&A",
        "감사합니다"
    )

    # 저장
    output_path = os.path.expanduser("~/sim2real/pen_grasp_sim2real_presentation.pptx")
    prs.save(output_path)
    print(f"프레젠테이션 저장 완료: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
